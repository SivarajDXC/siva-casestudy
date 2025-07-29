import os
import faiss
import numpy as np
import streamlit as st
from pathlib import Path
from io import BytesIO
from typing import Dict
from sentence_transformers import SentenceTransformer
from openai import AzureOpenAI
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from IPython.display import display, Markdown

class DocumentProcessor:
    def __init__(self):
        self.embedding_model = SentenceTransformer('all-mpnet-base-v2')
        self.index = None
        self.text_data = []

    def load_documents(self, uploaded_files):
        """Process uploaded documents"""
        for uploaded_file in uploaded_files:
            file_extension = Path(uploaded_file.name).suffix.lower()
            content = uploaded_file.getvalue()

            if file_extension == '.pdf':
                self._process_pdf(content)
            elif file_extension == '.pptx':
                self._process_pptx(content)
            elif file_extension == '.docx':
                self._process_docx(content)

    def _process_pdf(self, content):
        reader = PdfReader(BytesIO(content))
        for page in reader.pages:
            text = page.extract_text()
            if text:
                self.text_data.append(text)

    def _process_pptx(self, content):
        prs = Presentation(BytesIO(content))
        for slide in prs.slides:
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            if text:
                self.text_data.append("\n".join(text))

    def _process_docx(self, content):
        doc = Document(BytesIO(content))
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        if full_text:
            self.text_data.append("\n".join(full_text))

    def create_embeddings(self):
        """Create FAISS index from document embeddings"""
        if not self.text_data:
            raise ValueError("No text data to process")

        embeddings = self.embedding_model.encode(self.text_data)
        dimension = embeddings.shape[1]
        self.index = faiss.IndexFlatL2(dimension)
        self.index.add(embeddings.astype(np.float32))

class AnswerEvaluator:
    def __init__(self):
        # self.azure_client = AzureOpenAI(
        #     api_key=azure_config['api_key'],
        #     api_version="2024-08-01-preview",
        #     azure_endpoint=azure_config['endpoint']
        # )
        self.azure_client = AzureOpenAI(
            api_key="",
            api_version="",
            azure_endpoint=""
        )
        self.text_processor = None
        self.evaluation_results = []

    def _get_relevant_context(self, query: str, k: int = 3) -> list:
        """Retrieve relevant context from textbooks"""
        if not self.text_processor or not self.text_processor.index:
            raise ValueError("Text processor not initialized")

        query_embedding = self.text_processor.embedding_model.encode([query])
        distances, indices = self.text_processor.index.search(query_embedding.astype(np.float32), k)
        return [self.text_processor.text_data[i] for i in indices[0]]

    def extract_text_from_answer(self, content: bytes, file_extension: str) -> str:
        """Extract text from student's answer file"""
        try:
            if file_extension == '.pdf':
                return self._extract_from_pdf(content)
            elif file_extension == '.pptx':
                return self._extract_from_pptx(content)
            elif file_extension == '.docx':
                return self._extract_from_docx(content)
            else:
                raise ValueError("Unsupported file format")
        except Exception as e:
            raise RuntimeError(f"Error processing file: {str(e)}")

    def _extract_from_pdf(self, content):
        reader = PdfReader(BytesIO(content))
        return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

    def _extract_from_pptx(self, content):
        prs = Presentation(BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _extract_from_docx(self, content):
        doc = Document(BytesIO(content))
        return "\n".join([para.text for para in doc.paragraphs if para.text])

    def evaluate_answer(self, student_id: str, answer_text: str, criteria: str):
        """Evaluate a single student answer"""
        try:
            context = self._get_relevant_context(answer_text)
            context_str = "\n".join(context)

            prompt = f"""
            Evaluate the student's answer based on the following criteria: {criteria}

            Relevant textbook context: {context_str}

            Student's answer: {answer_text}

            For each answer against the questions mentioned in the answer text, provide a score
            from 0-10 based on the evaluation considering both the textbook material and
            general knowledge. Also provide a detailed explanation for the scoring and highlight
            strengths, weaknesses, and areas for improvement.
            """

            response = self.azure_client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3
            )

            self.evaluation_results.append({
                "Student ID": student_id,
                "Answer": answer_text,
                "Evaluation": response.choices[0].message.content,
                "Criteria": criteria
            })
        except Exception as e:
            raise RuntimeError(f"Evaluation failed: {str(e)}")

def main():
    # port = os.environ.get('PORT',8501)
    # st.run(port=int(port))

    
    # Streamlit UI
    st.set_page_config(page_title="Case Study Evaluation System", layout="wide")

    # Initialize session state
    if 'processor' not in st.session_state:
        st.session_state.processor = None
    if 'evaluations' not in st.session_state:
        st.session_state.evaluations = []

    # Main app
    st.title("Case Study Evaluation System")

    tab1, tab2 = st.tabs(["Train", "Evaluate"])

    # Section 1: Textbook Processing
    with tab1:
        st.header("1. Upload and Process Training Documents")
        uploaded_textbooks = st.file_uploader(
            "Upload textbook files (PDF, DOCX, PPTX)",
            type=["pdf", "docx", "pptx"],
            accept_multiple_files=True
        )

        if uploaded_textbooks and st.button("Train Model"):
            try:
                processor = DocumentProcessor()
                processor.load_documents(uploaded_textbooks)
                processor.create_embeddings()
                st.session_state.processor = processor
                st.success("Training complete! Documents are processed and embedded.")
            except Exception as e:
                st.error(f"Error processing textbooks: {str(e)}")

    # Section 2: Answer Evaluation
    with tab2:
        if st.session_state.processor:
            st.header("2. Evaluate Case Studies")

            uploaded_answer = st.file_uploader(
                "Upload student answer (PDF, DOCX, PPTX)",
                type=["pdf", "docx", "pptx"]
            )
            student_id = st.text_input("Student ID")
            criteria = st.text_area("Evaluation Criteria", "Accuracy of information")

            if uploaded_answer and student_id and criteria and st.button("Evaluate"):
                try:
                    evaluator = AnswerEvaluator()
                    evaluator.text_processor = st.session_state.processor

                    file_extension = Path(uploaded_answer.name).suffix.lower()
                    content = uploaded_answer.getvalue()
                    answer_text = evaluator.extract_text_from_answer(content, file_extension)

                    evaluator.evaluate_answer(student_id, answer_text, criteria)
                    st.session_state.evaluations.extend(evaluator.evaluation_results)
                    st.success("Evaluation completed successfully!")
                except Exception as e:
                    st.error(f"Evaluation error: {str(e)}")

    # Section 3: Display Results
    if st.session_state.evaluations:
        st.header("Evaluation Results")

        for idx, result in enumerate(st.session_state.evaluations):
            with st.expander(f"Evaluation for Student {result['Student ID']}"):
                st.subheader(f"Student ID: {result['Student ID']}")
                st.write(f"**Criteria:** {result['Criteria']}")

                col1, col2 = st.columns(2)
                with col1:
                    st.write("**Student Answer:**")
                    st.write(result['Answer'])
                with col2:
                    st.write("**Evaluation:**")
                    st.write(result['Evaluation'])

        # Download button
        results_text = []
        for result in st.session_state.evaluations:
            results_text.append(f"Student ID: {result['Student ID']}\n")
            results_text.append(f"Criteria: {result['Criteria']}\n")
            results_text.append("="*50 + "\n")
            results_text.append("Student Answer:\n")
            results_text.append(result['Answer'] + "\n\n")
            results_text.append("Evaluation:\n")
            results_text.append(result['Evaluation'] + "\n")
            results_text.append("="*50 + "\n\n")

        st.download_button(
            label="Download All Results",
            data="".join(results_text),
            file_name="evaluations.txt",
            mime="text/plain"
        )
    else:
        st.info("No evaluations yet. Process training documents and evaluate answers to see results.")

    # Section 4: System Status
    st.header("System Status")
    if st.session_state.processor:
        st.success("✅ Training documents index loaded and ready for evaluations")
        st.write(f"Number of textbook passages indexed: {len(st.session_state.processor.text_data)}")
    else:
        st.warning("⚠️ Please upload and process training documents first")

if __name__ == "__main__":
    main()
